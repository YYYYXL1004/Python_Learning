import os

from pptx import Presentation
from pptx.util import Inches


# PPTX文件保存路径
SAVE_DIR = os.path.join(os.path.dirname(__file__), "exp6_files")
PPTX_FILE = os.path.join(SAVE_DIR, "table_demo.pptx")


def add_table(slide, data):
    """在幻灯片中添加表格并写入数据。"""
    rows = len(data)
    cols = len(data[0])

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(0.8 + rows * 0.4)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for i in range(rows):
        for j in range(cols):
            table.cell(i, j).text = str(data[i][j])


def make_test_pptx():
    """准备一个包含表格的PPTX文件。"""
    os.makedirs(SAVE_DIR, exist_ok=True)

    prs = Presentation()
    layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "学生成绩表"
    data = [
        ["姓名", "语文", "数学", "英语"],
        ["张三", 86, 90, 88],
        ["李四", 78, 85, 82],
        ["王五", 92, 88, 95]
    ]
    add_table(slide, data)

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "课程信息表"
    data = [
        ["课程", "教师", "课时"],
        ["Python", "刘老师", 48],
        ["数据库", "王老师", 40],
        ["数据结构", "赵老师", 56]
    ]
    add_table(slide, data)

    prs.save(PPTX_FILE)


def read_table(table):
    """读取一个表格中的所有单元格内容。"""
    result = []

    for row in table.rows:
        line = []
        for cell in row.cells:
            line.append(cell.text)
        result.append(line)

    return result


def read_pptx_tables(file_path):
    """读取PPTX文件中所有表格的内容。"""
    prs = Presentation(file_path)
    table_count = 0

    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_table:
                table_count += 1
                print(f"第{slide_index}张幻灯片，第{table_count}个表格：")

                table_data = read_table(shape.table)
                for row in table_data:
                    print("\t".join(row))

                print()

    if table_count == 0:
        print("该PPTX文件中没有表格。")


def main():
    if not os.path.exists(PPTX_FILE):
        make_test_pptx()

    read_pptx_tables(PPTX_FILE)


if __name__ == "__main__":
    main()
