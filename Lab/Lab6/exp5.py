import os

from pptx import Presentation


# 保存本实验测试文件的文件夹
SAVE_DIR = os.path.join(os.path.dirname(__file__), "exp5_files")

# 使用全局变量统计PPTX文件数量和幻灯片总数量
pptx_count = 0
slide_count = 0


def make_pptx(file_path, count):
    """生成指定幻灯片数量的PPTX文件。"""
    prs = Presentation()
    layout = prs.slide_layouts[0]

    for i in range(1, count + 1):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"第{i}张幻灯片"
        slide.placeholders[1].text = "这是用于统计测试的幻灯片。"

    prs.save(file_path)


def make_test_files():
    """准备带有多个子文件夹的PPTX测试文件。"""
    folders = [
        SAVE_DIR,
        os.path.join(SAVE_DIR, "课件1"),
        os.path.join(SAVE_DIR, "课件2"),
        os.path.join(SAVE_DIR, "课件2", "子目录")
    ]

    for folder in folders:
        os.makedirs(folder, exist_ok=True)

    make_pptx(os.path.join(SAVE_DIR, "首页.pptx"), 2)
    make_pptx(os.path.join(SAVE_DIR, "课件1", "Python基础.pptx"), 3)
    make_pptx(os.path.join(SAVE_DIR, "课件2", "函数.pptx"), 4)
    make_pptx(os.path.join(SAVE_DIR, "课件2", "子目录", "递归.pptx"), 5)


def count_pptx_slides(folder):
    """递归遍历文件夹，统计所有PPTX文件中的幻灯片数量。"""
    global pptx_count, slide_count

    for name in os.listdir(folder):
        path = os.path.join(folder, name)

        if os.path.isdir(path):
            # 如果是文件夹，继续递归遍历
            count_pptx_slides(path)
        elif os.path.isfile(path) and name.lower().endswith(".pptx"):
            prs = Presentation(path)
            pptx_count += 1
            slide_count += len(prs.slides)
            print(f"{path}：{len(prs.slides)}张")


def main():
    global pptx_count, slide_count

    make_test_files()

    # 每次统计前先把全局变量清零
    pptx_count = 0
    slide_count = 0

    count_pptx_slides(SAVE_DIR)

    print(f"共统计PPTX文件{pptx_count}个")
    print(f"幻灯片总数量为{slide_count}张")


if __name__ == "__main__":
    main()
